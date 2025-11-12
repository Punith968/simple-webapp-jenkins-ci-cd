from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Helpers

ACCENT_BLUE = RGBColor(15, 108, 189)   # #0F6CBD
ACCENT_GREEN = RGBColor(46, 182, 125)  # #2EB67D
BG_LIGHT = RGBColor(248, 250, 252)
TEXT_DARK = RGBColor(30, 41, 59)


def style_slide_background(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text="Punith C  |  Nov 13, 2025"):
    left, top, width, height = Inches(0.5), Inches(6.8), Inches(9), Inches(0.3)
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 116, 139)
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    style_slide_background(slide)
    # Title
    slide.shapes.title.text = title
    ttf = slide.shapes.title.text_frame.paragraphs[0].font
    ttf.size = Pt(38)
    ttf.bold = True
    ttf.color.rgb = TEXT_DARK
    # Subtitle
    slide.placeholders[1].text = subtitle
    stf = slide.placeholders[1].text_frame.paragraphs[0].font
    stf.size = Pt(18)
    stf.color.rgb = ACCENT_BLUE
    add_footer(slide)
    return slide


def add_bullet_slide(prs, title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    style_slide_background(slide)
    slide.shapes.title.text = title
    # Title style
    title_font = slide.shapes.title.text_frame.paragraphs[0].font
    title_font.size = Pt(28)
    title_font.bold = True
    title_font.color.rgb = ACCENT_BLUE
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_footer(slide)
    return slide


def add_section_header(prs, title):
    slide_layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(slide_layout)
    style_slide_background(slide)
    slide.shapes.title.text = title
    title_font = slide.shapes.title.text_frame.paragraphs[0].font
    title_font.size = Pt(32)
    title_font.bold = True
    title_font.color.rgb = ACCENT_GREEN
    add_footer(slide)
    return slide


def add_placeholder_image_slide(prs, title, caption):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    style_slide_background(slide)
    slide.shapes.title.text = title
    title_font = slide.shapes.title.text_frame.paragraphs[0].font
    title_font.size = Pt(28)
    title_font.bold = True
    title_font.color.rgb = ACCENT_BLUE
    # Draw a placeholder rectangle for where to put screenshots later
    left, top, width, height = Inches(1), Inches(1.8), Inches(8), Inches(4)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(230, 230, 230)
    line = shape.line
    line.color.rgb = RGBColor(120, 120, 120)
    tx_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(1))
    tf = tx_box.text_frame
    tf.text = caption
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_footer(slide)
    return slide


def add_image_slide(prs, title, image_path, placeholder_caption):
    if image_path and os.path.exists(image_path):
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)
        style_slide_background(slide)
        slide.shapes.title.text = title
        title_font = slide.shapes.title.text_frame.paragraphs[0].font
        title_font.size = Pt(28)
        title_font.bold = True
        title_font.color.rgb = ACCENT_BLUE
        left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
        slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        add_footer(slide)
        return slide
    else:
        return add_placeholder_image_slide(prs, title, placeholder_caption)


def main():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "Jenkins CI/CD for a Static Web App — Case Study",
        "End-to-end CI/CD with Docker deployment | Punith C | Nov 13, 2025",
    )

    # Slide 2: Executive Summary
    add_bullet_slide(
        prs,
        "Executive Summary",
        [
            "Goal: Automate build → test → deploy for a static site",
            "Pipeline: Checkout, Build, Test, Deploy (4/4 green)",
            "Deployment: Docker (nginx) on port 8090",
            "Result: 100% success; ~15s from commit to running app",
        ],
        notes="Production-ready automation with post-deploy verification.",
    )

    # Slide 3: Problem Statement & Objectives
    add_bullet_slide(
        prs,
        "Problem Statement & Objectives",
        [
            "Build a CI/CD pipeline for a static website (HTML/CSS/JS)",
            "Automated validation and reliable deployment",
            "Use Jenkins Declarative Pipeline",
            "Prefer Docker for reproducibility",
        ],
    )

    # Slide 4: Architecture Overview (diagram placeholder)
    add_bullet_slide(
        prs,
        "Architecture Overview",
        [
            "Flow: Developer (GitHub) → Jenkins → Docker Image → Nginx Container → Browser",
            "Environment: Windows host + WSL2 Ubuntu agent",
            "Base image: nginx:alpine",
        ],
        notes="Add a diagram later with boxes and arrows.",
    )

    # Slide 5: Technologies Used
    add_bullet_slide(
        prs,
        "Technologies Used",
        ["HTML5, CSS3, JavaScript", "Nginx (nginx:alpine)", "Docker", "Jenkins (Declarative)", "Git/GitHub"],
    )

    # Slide 6: Tools & Environment
    add_bullet_slide(
        prs,
        "Tools & Environment",
        [
            "Jenkins 2.528.2, Docker 28.2.2, Git 2.43.0",
            "Windows + WSL2 (Ubuntu 24.04)",
            "VS Code; PowerShell + bash",
            "Optional: docker-compose",
        ],
        notes="Jenkins user added to docker group; restarted service.",
    )

    # Slide 7: Repository & Files
    add_bullet_slide(
        prs,
        "Repository & Files",
        [
            "App: index.html, styles.css, script.js",
            "Pipeline: Jenkinsfile (Declarative)",
            "Container: Dockerfile, docker-compose.yml",
            "Scripts: validate.sh, deploy.sh (alt method)",
            "Docs: README, CASE-STUDY-SUMMARY, DEPLOYMENT-GUIDE",
        ],
    )

    # Slide 8: Pipeline Design
    add_bullet_slide(
        prs,
        "Pipeline Design (Stages)",
        [
            "Checkout → Build (validate.sh) → Test (validate.sh) → Deploy (docker build/run)",
            "Env: DEPLOY_METHOD=docker, CONTAINER_NAME=simple-webapp-demo, WEBAPP_PORT=8090",
            "dir() wrappers handle nested repo path",
        ],
    )

    # Slide 9: Implementation Steps
    add_bullet_slide(
        prs,
        "Implementation Steps",
        [
            "Setup repo, scripts; push to GitHub",
            "Configure Jenkins (Pipeline from SCM; correct script path)",
            "chmod +x scripts; ensure Docker access for Jenkins",
            "Pre-pull nginx:alpine to avoid timeouts",
        ],
    )

    # Slide 10: Deployment Approach (Docker)
    add_bullet_slide(
        prs,
        "Deployment Approach (Docker)",
        [
            "Build image: simple-webapp:latest",
            "Container: simple-webapp-demo (8090→80)",
            "Stop/remove old container automatically",
            "Post actions verify container health and URL",
        ],
    )

    # Slide 11: Results & Metrics
    add_bullet_slide(
        prs,
        "Results & Metrics",
        [
            "Build #10: SUCCESS (4/4 stages), ~15s total",
            "Image ID: e54b76dbdf5f; Container ID: 1db05a33b511",
            "URL: http://localhost:8090",
            "JS alert verified: ‘Hello from the Simple WebApp Jenkins CI/CD demo!’",
        ],
    )

    # Slide 12: Challenges & Resolutions
    add_bullet_slide(
        prs,
        "Challenges & Resolutions",
        [
            "Jenkinsfile path in nested dir → set Script Path + use dir()",
            "Script permissions → git chmod + pipeline chmod",
            "Docker permission denied → add Jenkins to docker group + restart",
            "Image pull timeout → pre-cache nginx:alpine",
            "Deploy method → switched from sudo copy to Docker",
        ],
    )

    # Slide 13: Lessons Learned
    add_bullet_slide(
        prs,
        "Lessons Learned",
        [
            "Docker simplifies deployments in CI/CD",
            "Repo structure impacts Jenkins configuration",
            "Windows→Linux: scripts need executable bits",
            "Group membership changes require service restart",
        ],
    )

    # Slide 14: Next Steps
    add_bullet_slide(
        prs,
        "Next Steps (Enhancements)",
        [
            "UI tests with Playwright/Puppeteer",
            "Push images to registry; tag by build number",
            "Blue/green or canary deployments",
            "Webhook triggers; Slack/Email notifications",
        ],
    )

    # Slide 15: Screenshots (embed if available, else placeholders)
    add_image_slide(
        prs,
        "Screenshot: Jenkins Build #10",
        os.path.join("screenshots", "jenkins_build.png"),
        "Place Jenkins build screenshot here (4/4 green stages)",
    )
    add_image_slide(
        prs,
        "Screenshot: Docker Container (8090→80)",
        os.path.join("screenshots", "docker_ps.png"),
        "Place 'docker ps' screenshot showing simple-webapp-demo here",
    )
    add_image_slide(
        prs,
        "Screenshot: App at http://localhost:8090",
        os.path.join("screenshots", "app_browser.png"),
        "Place browser screenshot with JS alert here",
    )

    # Slide 16: Thank You
    add_bullet_slide(
        prs,
        "Thank You",
        [
            "Repo: https://github.com/Punith968/simple-webapp-jenkins-ci-cd",
            "Questions welcome",
        ],
    )

    prs.save("project_presentation.pptx")
    print("Saved: project_presentation.pptx")


if __name__ == "__main__":
    main()
