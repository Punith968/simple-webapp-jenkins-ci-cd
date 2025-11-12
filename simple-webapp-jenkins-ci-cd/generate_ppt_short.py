from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

ACCENT_BLUE = RGBColor(15, 108, 189)
BG_LIGHT = RGBColor(248, 250, 252)
TEXT_DARK = RGBColor(30, 41, 59)


def style_slide(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_LIGHT
    # brand bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.25))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_BLUE; bar.line.fill.background()


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def title_slide(prs, title, subtitle, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    style_slide(s)
    s.shapes.title.text = title
    s.placeholders[1].text = subtitle
    if notes: add_notes(s, notes)


def bullets(prs, title, items, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    style_slide(s)
    s.shapes.title.text = title
    tf = s.shapes.placeholders[1].text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it; p.level = 0; p.font.size = Pt(20)
    if notes: add_notes(s, notes)


def image_or_placeholder(prs, title, path, caption, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    style_slide(s)
    s.shapes.title.text = title
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    if os.path.exists(path):
        s.shapes.add_picture(path, left, top, width=width, height=height)
    else:
        ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(230,230,230); ph.line.color.rgb = RGBColor(120,120,120)
        tb = s.shapes.add_textbox(left+Inches(0.2), top+Inches(0.2), width-Inches(0.4), Inches(0.8))
        tb.text_frame.text = caption
    if notes: add_notes(s, notes)


def main():
    prs = Presentation()

    title_slide(prs, "Jenkins CI/CD — Static Web App", "Short Deck | Punith C | Nov 13, 2025",
                notes="Quick overview: goals, pipeline, results.")

    bullets(prs, "Summary", [
        "Goal: CI/CD for static site",
        "Pipeline: Checkout, Build, Test, Deploy",
        "Deployment: Docker (nginx) on 8090",
        "Result: 100% success; ~15s",
    ], notes="Emphasize production-ready demo.")

    bullets(prs, "Architecture", [
        "GitHub → Jenkins → Docker Image → Nginx Container → Browser",
        "Windows host; WSL2 agent",
        "Base: nginx:alpine",
    ], notes="Keep it high-level.")

    bullets(prs, "Pipeline (Declarative)", [
        "Checkout → Build (validate.sh) → Test (validate.sh) → Deploy (docker)",
        "Env: DEPLOY_METHOD=docker; PORT=8090",
    ], notes="Call out dir() usage for nested path.")

    bullets(prs, "Implementation", [
        "Make scripts executable; configure Jenkins",
        "Add Jenkins to docker group; restart",
        "Pre-pull nginx:alpine",
    ], notes="Key setup steps to avoid failures.")

    bullets(prs, "Deployment", [
        "Image: simple-webapp:latest",
        "Container: simple-webapp-demo (8090→80)",
        "Auto-stop old container on redeploy",
    ], notes="Why Docker over sudo copy.")

    bullets(prs, "Results", [
        "Build #10: SUCCESS (4/4)",
        "~15s total runtime",
        "URL: http://localhost:8090",
        "JS alert verified",
    ], notes="Show concrete proof.")

    image_or_placeholder(prs, "Jenkins Build", os.path.join('screenshots','jenkins_build.png'), "Add Jenkins 4/4 green screenshot",
                         notes="Replace with actual screenshot.")

    bullets(prs, "Lessons & Next", [
        "Docker simplifies CI/CD",
        "Repo structure matters (dir(), script path)",
        "Next: UI tests, registry, blue/green",
    ], notes="Close with what’s next.")

    bullets(prs, "Thank You", [
        "Repo: github.com/Punith968/simple-webapp-jenkins-ci-cd",
        "Questions welcome",
    ])

    prs.save('project_presentation_short.pptx')
    print('Saved: project_presentation_short.pptx')


if __name__ == '__main__':
    main()
